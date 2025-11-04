"""PowerPoint file parser."""

from pathlib import Path
from .base import BaseParser, ParseResult


class PPTXParser(BaseParser):
    """PowerPoint 파일 파서 (python-pptx 사용)"""

    def parse(self, file_path: Path) -> ParseResult:
        """
        PPTX 파일에서 텍스트를 추출합니다.

        Args:
            file_path: 파싱할 PPTX 파일 경로

        Returns:
            ParseResult: 파싱된 텍스트와 메타데이터
        """
        if not file_path.exists():
            raise FileNotFoundError(f"파일을 찾을 수 없습니다: {file_path}")

        try:
            try:
                from pptx import Presentation
            except ImportError:
                raise ImportError(
                    "python-pptx가 설치되지 않았습니다. 'pip install python-pptx'로 설치하세요."
                )

            prs = Presentation(file_path)

            # 슬라이드별로 텍스트 추출
            slides_text = []
            for idx, slide in enumerate(prs.slides, 1):
                slide_content = [f"=== Slide {idx} ==="]

                # 슬라이드의 모든 도형에서 텍스트 추출
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())

                    # 테이블이 있는 경우 처리
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                            if row_text.strip():
                                slide_content.append(row_text)

                # 슬라이드 노트 추출
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_content.append(f"Notes: {notes_text}")

                slides_text.append('\n'.join(slide_content))

            content = '\n\n'.join(slides_text)

            metadata = {
                'file_type': 'pptx',
                'num_slides': len(prs.slides),
                'file_size': file_path.stat().st_size,
            }

            # 슬라이드 크기 정보
            if prs.slide_width and prs.slide_height:
                metadata['slide_width'] = prs.slide_width
                metadata['slide_height'] = prs.slide_height

            return ParseResult(content=content, metadata=metadata)

        except Exception as e:
            raise ValueError(f"PPTX 파싱 실패: {file_path}") from e

    def supports(self, file_path: Path) -> bool:
        """PPTX 파일 확장자 지원 확인"""
        return file_path.suffix.lower() in self.supported_extensions

    @property
    def supported_extensions(self) -> list[str]:
        """지원하는 PowerPoint 파일 확장자"""
        return ['.pptx']
